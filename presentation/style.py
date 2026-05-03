"""Slide style definitions and helpers."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.util import Inches, Pt

# ── Colour palette ────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)   # near-black
BG_CARD    = RGBColor(0x16, 0x1E, 0x2B)   # dark navy
ACCENT1    = RGBColor(0x00, 0xE5, 0xFF)   # cyan
ACCENT2    = RGBColor(0x76, 0xFF, 0x03)   # green
ACCENT3    = RGBColor(0xFF, 0x6D, 0x00)   # orange
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID   = RGBColor(0x8A, 0x9B, 0xAE)
RED_ERR    = RGBColor(0xFF, 0x3D, 0x00)
NVIDIA_GRN = RGBColor(0x76, 0xB9, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def fill_bg(slide, color=BG_DARK):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=BG_CARD, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=Pt(14), bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    from pptx.util import Pt
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_title_bar(slide, title, subtitle=None):
    """Dark gradient title bar at top."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=BG_CARD)
    # Left accent stripe
    add_rect(slide, 0, 0, Inches(0.08), Inches(1.2), fill_color=ACCENT1)

    add_textbox(slide, Inches(0.25), Inches(0.1), Inches(12), Inches(0.65),
                title, font_size=Pt(28), bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.25), Inches(0.75), Inches(12), Inches(0.4),
                    subtitle, font_size=Pt(13), color=ACCENT1)


def add_section_label(slide, label, x, y, w=Inches(3)):
    """Small coloured section label."""
    add_textbox(slide, x, y, w, Inches(0.3), label.upper(),
                font_size=Pt(9), bold=True, color=ACCENT1)


def speedup_color(speedup):
    """Color depending on speedup bucket."""
    if speedup is None:
        return RED_ERR
    if speedup >= 5:
        return ACCENT2
    if speedup >= 3:
        return ACCENT1
    if speedup >= 2:
        return WHITE
    return GREY_MID


def bullet_list(slide, items, x, y, w, h, font_size=Pt(13), color=WHITE,
                bullet_color=ACCENT1, indent=Inches(0.25)):
    """Add a bulleted list of strings."""
    for i, item in enumerate(items):
        row_y = y + i * (font_size + Pt(4))
        if row_y + font_size > y + h:
            break
        # bullet dot
        add_textbox(slide, x, row_y, Inches(0.2), Inches(0.3), "▸",
                    font_size=font_size, color=bullet_color)
        add_textbox(slide, x + indent, row_y, w - indent, Inches(0.35), item,
                    font_size=font_size, color=color, wrap=True)
